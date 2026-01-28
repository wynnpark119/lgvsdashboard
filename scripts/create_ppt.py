#!/usr/bin/env python3
"""LG VS MTK Dashboard Overview PPT Generator"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# LG 브랜드 컬러
LG_RED = RGBColor(165, 0, 52)
DARK_GRAY = RGBColor(51, 51, 51)
LIGHT_GRAY = RGBColor(245, 245, 245)

def add_title_slide(prs, title, subtitle=""):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Background shape
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), Inches(10), Inches(2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LG_RED
    shape.line.fill.background()
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.7), Inches(9), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(44)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Subtitle
    if subtitle:
        txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.5), Inches(9), Inches(0.5))
        tf2 = txBox2.text_frame
        p2 = tf2.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(20)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
    
    return slide

def add_content_slide(prs, title, bullets):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LG_RED
    shape.line.fill.background()
    
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Content
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(5))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    
    for i, bullet in enumerate(bullets):
        if i == 0:
            p = tf2.paragraphs[0]
        else:
            p = tf2.add_paragraph()
        p.text = f"• {bullet}"
        p.font.size = Pt(20)
        p.font.color.rgb = DARK_GRAY
        p.space_after = Pt(12)
    
    return slide

def add_table_slide(prs, title, headers, rows):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Title bar
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(10), Inches(1.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LG_RED
    shape.line.fill.background()
    
    # Title text
    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.6))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    
    # Table
    num_cols = len(headers)
    num_rows = len(rows) + 1
    
    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.5), Inches(1.5), Inches(9), Inches(0.5 * num_rows)).table
    
    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = DARK_GRAY
        for paragraph in cell.text_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.bold = True
            paragraph.font.color.rgb = RGBColor(255, 255, 255)
    
    # Data rows
    for row_idx, row in enumerate(rows):
        for col_idx, cell_text in enumerate(row):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = cell_text
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.size = Pt(12)
                paragraph.font.color.rgb = DARK_GRAY
    
    return slide

def add_section_slide(prs, section_num, title):
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Number circle
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(4.25), Inches(2), Inches(1.5), Inches(1.5))
    shape.fill.solid()
    shape.fill.fore_color.rgb = LG_RED
    shape.line.fill.background()
    
    # Number
    txBox = slide.shapes.add_textbox(Inches(4.25), Inches(2.3), Inches(1.5), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = str(section_num)
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER
    
    # Title
    txBox2 = slide.shapes.add_textbox(Inches(0.5), Inches(3.8), Inches(9), Inches(1))
    tf2 = txBox2.text_frame
    p2 = tf2.paragraphs[0]
    p2.text = title
    p2.font.size = Pt(36)
    p2.font.bold = True
    p2.font.color.rgb = DARK_GRAY
    p2.alignment = PP_ALIGN.CENTER
    
    return slide

def main():
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)
    
    # Slide 1: Title
    add_title_slide(prs, "LG VS MTK Dashboard", "기술 관심 모니터링 대시보드 개요")
    
    # Slide 2: 대시보드 목적
    add_content_slide(prs, "대시보드 목적", [
        "OEM의 LG VS 기술에 대한 관심 상태 파악",
        '"이해를 동반한 접촉이 발생한 기술" 식별',
        "기술 조직과의 협업 시점 판단",
        '"많이 본 기술"이 아니라 "깊이 있게 검토된 기술" 강조'
    ])
    
    # Slide 3: 전체 구조
    add_section_slide(prs, 1, "대시보드 전체 구조")
    
    # Slide 4: Page Architecture
    add_table_slide(prs, "Page Architecture", 
        ["섹션", "경로", "설명"],
        [
            ["전체 요약", "/", "B2B 퍼널 현황, 기술별 관심 신호"],
            ["기술별 관심 현황", "/review-flow/*", "TOFU/MOFU/BOFU 단계별 현황"],
            ["캠페인/광고 영향", "/campaign-impact/*", "캠페인 전·후 관심 단계 이동"],
            ["상세 데이터", "/detail/*", "기술별/채널별/문의/콘텐츠 상세"],
        ]
    )
    
    # Slide 5: 기술별 관심 현황 상세
    add_table_slide(prs, "기술별 관심 현황 (Progress)", 
        ["페이지", "경로", "설명"],
        [
            ["전체 현황", "/review-flow", "기술별 관심 분포 및 추이"],
            ["처음 접촉 (TOFU)", "/review-flow/tofu", "최초 유입, 인지 형성 여부"],
            ["깊은 관심 (MOFU/BOFU)", "/review-flow/mofu-bofu", "반복 접촉, 심화 콘텐츠 소비"],
            ["기술별 비교", "/review-flow/technology", "기술 간 관심도 비교"],
            ["추세", "/review-flow/momentum", "시간에 따른 관심 변화"],
        ]
    )
    
    # Slide 6: B2B 퍼널
    add_section_slide(prs, 2, "B2B 퍼널 정의")
    
    # Slide 7: 퍼널 정의 테이블
    add_table_slide(prs, "B2B 퍼널 단계", 
        ["단계", "정의", "주요 신호"],
        [
            ["TOFU (처음 접촉)", "최초 유입, 인지 형성", "첫 방문, 단일 콘텐츠 조회"],
            ["MOFU (반복 접촉)", "관심 지속, 심화 탐색", "재방문, 복수 콘텐츠, 웨비나"],
            ["BOFU (깊은 관심)", "구체적 관심 표현", "Inquiry, 스펙 요청, 다운로드"],
        ]
    )
    
    # Slide 8: 데이터 소스
    add_section_slide(prs, 3, "주요 데이터 소스")
    
    # Slide 9: 데이터 소스 테이블
    add_table_slide(prs, "현행 수집 범위", 
        ["소스", "수집 항목"],
        [
            ["LG.com", "기술/캠페인 페이지 방문, 체류 시간, 다운로드"],
            ["LinkedIn", "기술 콘텐츠 클릭, Engagement (Organic/Paid)"],
            ["YouTube", "기술 영상 조회, 시청 깊이 (%)"],
            ["Inquiry", "문의 폼 제출, 스펙 요청"],
            ["Newsletter", "구독자 수, 오픈율, 클릭율"],
        ]
    )
    
    # Slide 10: 성과 지표
    add_section_slide(prs, 4, "대표 성과 지표")
    
    # Slide 11: 핵심 지표
    add_content_slide(prs, "핵심 지표 (Primary)", [
        "TOFU → MOFU 전환율: 재방문 발생 비율",
        "MOFU → BOFU 전환율: Inquiry/스펙 요청 발생 비율",
        "기술별 관심 추세: 30/60/90일 관심도 변화",
        "콘텐츠 완료율: 끝까지 소비한 비율 (70%+ = 높은 관심)",
        "뉴스레터 구독자 순증: 월간 신규 - 해지"
    ])
    
    # Slide 12: 판단 기준
    add_table_slide(prs, "관심 판단 기준", 
        ["신호", "의미"],
        [
            ["체류 시간 3분+", "콘텐츠 깊이 있게 소비"],
            ["영상 50%+ 시청", "기술에 대한 실질적 관심"],
            ["복수 접점 발생", "단일 기술에 대한 지속적 관심"],
            ["Inquiry/스펙 요청", "구체적 검토 의사 표현"],
            ["뉴스레터 구독", "지속적 관심 유지 의향"],
        ]
    )
    
    # Slide 13: 대시보드 역할
    add_section_slide(prs, 5, "대시보드의 역할")
    
    # Slide 14: 역할 상세
    add_content_slide(prs, "이 대시보드가 하는 것", [
        "기술 관심의 현재 상태 파악 (증가/유지/하락)",
        "관심 단계 식별 (처음 접촉 vs 깊은 관심)",
        "기술 조직 협업 시점 판단 (BOFU 진입 시)",
        "캠페인 효과 해석 (Amplifier/Accelerator/Noise)"
    ])
    
    # Slide 15: 하지 않는 것
    add_content_slide(prs, "이 대시보드가 하지 않는 것", [
        "CTR/CPC/ROI 평가 → 성과 마케팅 KPI는 별도 도구",
        "영업 리드 전달 → LG VS 마케팅팀 역할 범위 아님",
        "자동 인사이트/추천 → 판단은 사용자 몫",
        "채널 간 효율 비교 → 채널별 역할이 다름"
    ])
    
    # Slide 16: 광고의 역할
    add_content_slide(prs, "광고(Paid Media)의 역할", [
        "독립 KPI로 평가하지 않음",
        '"관심 흐름을 얼마나 증폭/가속했는가"로 해석',
        "광고 종료 후에도 관심이 유지되는지 확인",
        "Amplifier / Accelerator / Noise로 분류"
    ])
    
    # Slide 17: 배포 정보
    add_section_slide(prs, 6, "배포 정보")
    
    # Slide 18: 기술 스택 & URL
    add_content_slide(prs, "기술 스택 & 배포", [
        "Frontend: Next.js 14, TypeScript, Tailwind CSS, Recharts",
        "Deployment: Cloudflare Pages",
        "Production URL: https://lgvs-dashboard.pages.dev",
        "Repository: github.com/wynnpark119/lgvsdashboard"
    ])
    
    # Save
    output_path = "/Users/wynn.park/Documents/LG_VS/docs/LG_VS_Dashboard_Overview.pptx"
    prs.save(output_path)
    print(f"PPT saved to: {output_path}")

if __name__ == "__main__":
    main()
